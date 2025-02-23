from typing import List, Dict, Tuple
import logging
from pathlib import Path
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_core.documents import Document
from .base import BaseDocumentLoader
from processors.image_extractor import ImageExtractor

logger = logging.getLogger(__name__)

class PPTLoader(BaseDocumentLoader):
    """PowerPoint document loader - extracts text and images from slides"""
    
    def __init__(self, file_path: str):
        """Initialize loader"""
        super().__init__(file_path)
        self.image_extractor = ImageExtractor()
        self.image_map = {}  # Map to store image positions
    
    def _extract_images(self, prs: Presentation) -> Dict[str, dict]:
        """Extract images from presentation and map them to their relationships"""
        images = {}
        image_index = 1
        
        # Process all slides
        for slide_idx, slide in enumerate(prs.slides, 1):  # 使用 enumerate 获取索引
            for rel in slide.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        # Get image data
                        image_data = rel.target_part.blob
                        img_stream = io.BytesIO(image_data)
                        
                        # Process image
                        logger.info(f"Processing embedded image {image_index} from slide {slide_idx}")
                        extraction_result = self.image_extractor.extract_info(img_stream)
                        
                        if extraction_result["status"] == "success":
                            images[rel.rId] = {
                                "index": image_index,
                                "slide": slide_idx,  # 使用 enumerate 的索引
                                "content": extraction_result["content"],
                                "status": "success"
                            }
                        else:
                            images[rel.rId] = {
                                "index": image_index,
                                "slide": slide_idx,  # 使用 enumerate 的索引
                                "error": extraction_result.get("error", "Unknown error"),
                                "status": "failed"
                            }
                        
                        image_index += 1
                        
                    except Exception as e:
                        logger.error(f"Failed to process image {image_index} from slide {slide_idx}: {str(e)}")
                        images[rel.rId] = {
                            "index": image_index,
                            "slide": slide_idx,  # 使用 enumerate 的索引
                            "error": str(e),
                            "status": "failed"
                        }
                        image_index += 1
        
        return images
    
    def _process_shape(self, shape) -> str:
        """Process a shape and extract its content"""
        content_parts = []
        
        # Extract text if shape has text
        if hasattr(shape, "text") and shape.text.strip():
            content_parts.append(shape.text.strip())
        
        # Handle different shape types
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                # Get image relationship ID from shape's XML
                blip_rId = None
                for element in shape._element.iter():
                    if element.tag.endswith('blip'):
                        blip_rId = element.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        break
                
                if blip_rId and blip_rId in self.image_map:
                    img = self.image_map[blip_rId]
                    if img["status"] == "success":
                        content_parts.append(
                            f'\n<image id="{img["index"]:03d}">\n'
                            f'{img["content"]}\n'
                            f'</image>\n'
                        )
                    else:
                        content_parts.append(
                            f'\n<image id="{img["index"]:03d}" status="failed">\n'
                            f'Image processing failed: {img.get("error", "Unknown error")}\n'
                            f'</image>\n'
                        )
            except Exception as e:
                logger.error(f"Failed to process picture shape: {str(e)}")
        
        # Process group shapes recursively
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                child_content = self._process_shape(child)
                if child_content:
                    content_parts.append(child_content)
        
        # Process tables
        elif hasattr(shape, "table"):
            rows = []
            for row in shape.table.rows:
                cells = []
                for cell in row.cells:
                    if cell.text.strip():
                        cells.append(cell.text.strip())
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                content_parts.append("\n".join(rows))
        
        return "\n".join(content_parts)
    
    def _process_slide(self, slide, slide_number: int) -> str:
        """Process a slide and extract its content"""
        content_parts = [f"\n=== Slide {slide_number} ===\n"]
        
        # Process slide title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
            if title:
                content_parts.append(f"Title: {title}\n")
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            content = self._process_shape(shape)
            if content:
                content_parts.append(content)
        
        return "\n".join(content_parts)
    
    def _extract_content(self, prs: Presentation) -> Tuple[str, List[dict]]:
        """Extract content from presentation"""
        content_parts = []
        self.image_map = self._extract_images(prs)
        
        # Process all slides
        for idx, slide in enumerate(prs.slides, 1):
            slide_content = self._process_slide(slide, idx)
            if slide_content:
                content_parts.append(slide_content)
        
        # Convert image map to list for metadata
        images = [img for img in self.image_map.values()]
        return "\n\n".join(content_parts), images
    
    def load(self) -> List[Document]:
        """Load PowerPoint document and extract content"""
        try:
            # Check file exists
            ppt_path = Path(self.file_path)
            if not ppt_path.exists():
                raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")
            
            # Load presentation
            logger.info(f"Loading PowerPoint document: {ppt_path}")
            prs = Presentation(ppt_path)
            
            # Extract content
            content, images = self._extract_content(prs)
            
            # Create Document object
            doc = Document(
                page_content=content,
                metadata={
                    "source": str(ppt_path),
                    "file_type": "powerpoint",
                    "file_name": ppt_path.name,
                    "total_slides": len(prs.slides),
                    "images": images
                }
            )
            
            return [doc]
            
        except Exception as e:
            logger.error(f"Error loading PowerPoint document: {str(e)}", exc_info=True)
            # Return error document
            error_doc = Document(
                page_content=f"Failed to load PowerPoint document: {str(e)}",
                metadata={
                    "source": str(self.file_path),
                    "file_type": "powerpoint",
                    "extraction_status": "failed",
                    "error": str(e)
                }
            )
            return [error_doc] 